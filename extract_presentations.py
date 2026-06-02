#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""Extract and summarize PowerPoint presentation content from semanas 8-11"""

from pptx import Presentation
import os
from pathlib import Path

presentations = [
    r'c:\pyProjects\PFACD\apresentacoes semanais\Apresentações Semanais\semana8_grupo22.pptx',
    r'c:\pyProjects\PFACD\apresentacoes semanais\Apresentações Semanais\semana9_grupo22.pptx',
    r'c:\pyProjects\PFACD\apresentacoes semanais\Apresentações Semanais\semana10_grupo22.pptx',
    r'c:\pyProjects\PFACD\apresentacoes semanais\Apresentações Semanais\semana11_grupo22.pptx'
]

summary_data = {}

for prs_path in presentations:
    if os.path.exists(prs_path):
        try:
            prs = Presentation(prs_path)
            filename = os.path.basename(prs_path)
            week_name = filename.replace('_grupo22.pptx', '')
            
            print(f'\n{"="*80}')
            print(f'{week_name.upper()}')
            print(f'{"="*80}')
            print(f'Ficheiro: {filename}')
            print(f'Total de slides: {len(prs.slides)}\n')
            
            week_data = {
                'week': week_name,
                'file': filename,
                'total_slides': len(prs.slides),
                'slides': []
            }
            
            for idx, slide in enumerate(prs.slides, 1):
                slide_content = {
                    'number': idx,
                    'texts': []
                }
                
                print(f'[Slide {idx}]')
                
                # Extract all text from shapes
                slide_texts = []
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text.strip():
                        slide_texts.append(shape.text.strip())
                
                if slide_texts:
                    for text in slide_texts:
                        # Clean up text and print
                        lines = text.split('\n')
                        for line in lines:
                            clean_line = line.strip()
                            if clean_line:
                                print(f'  • {clean_line}')
                                slide_content['texts'].append(clean_line)
                else:
                    print('  (Slide sem conteúdo de texto)')
                
                print()
                week_data['slides'].append(slide_content)
            
            summary_data[week_name] = week_data
            
        except Exception as e:
            print(f'Erro ao processar {prs_path}: {str(e)}')
    else:
        print(f'Ficheiro não encontrado: {prs_path}')

# Print summary
print(f'\n{"="*80}')
print('RESUMO POR SEMANA')
print(f'{"="*80}\n')

for week in ['semana8', 'semana9', 'semana10', 'semana11']:
    if week in summary_data:
        data = summary_data[week]
        print(f'{week.upper()}:')
        print(f'  - Slides: {data["total_slides"]}')
        main_topics = []
        for slide in data['slides']:
            if slide['texts']:
                main_topics.extend([t[:50] + '...' if len(t) > 50 else t for t in slide['texts'][:2]])
        
        print(f'  - Principais tópicos:')
        for i, topic in enumerate(main_topics[:5], 1):
            print(f'    {i}. {topic}')
        print()
